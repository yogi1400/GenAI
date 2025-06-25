import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm
from io import BytesIO
import re
from fastapi import APIRouter, HTTPException, Request
from fastapi.responses import FileResponse
from tempfile import NamedTemporaryFile
import logging

# --- OpenRouter Kimi (moonshotai/kimi-dev-72b:free) Integration ---
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY", "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")
OPENROUTER_API_URL = "https://openrouter.ai/api/v1/chat/completions"

# --- Pexels API ---
PEXELS_API_KEY = os.getenv("PEXELS_API_KEY", "6N6GGbdSM0XT7oMfcStZTylIQpEvchwqFtCny7fGBlWeQpPyaxkISwaF")

router = APIRouter()

# Add logging for debugging
logger = logging.getLogger("presentation_creator")
logging.basicConfig(level=logging.INFO)

def openrouter_chat(messages, model="moonshotai/kimi-dev-72b:free", max_tokens=512):
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": model,
        "messages": messages,
        "max_tokens": max_tokens
    }
    logger.info(f"[DEBUG] Sending request to OpenRouter: {payload}")
    response = requests.post(OPENROUTER_API_URL, headers=headers, json=payload, timeout=60)
    logger.info(f"[DEBUG] OpenRouter response status: {response.status_code}")
    result = response.json()
    logger.info(f"[DEBUG] OpenRouter response: {result}")
    if response.status_code != 200 or "error" in result:
        error_msg = result.get("error", {}).get("message", response.text)
        logger.error(f"[ERROR] OpenRouter API error: {error_msg}")
        raise RuntimeError(f"OpenRouter API error: {error_msg}")
    if "choices" not in result or not result["choices"]:
        logger.error(f"[ERROR] OpenRouter API returned no choices: {result}")
        raise RuntimeError(f"OpenRouter API returned no choices: {result}")
    content = result["choices"][0]["message"]["content"]
    content = re.sub(r'<think>[\s\S]*?<\/think>', '', content, flags=re.IGNORECASE)
    return content.strip()

def search_image(query):
    headers = {"Authorization": PEXELS_API_KEY}
    url = f"https://api.pexels.com/v1/search?query={query}&per_page=1"
    response = requests.get(url, headers=headers)
    data = response.json()
    if data.get("photos"):
        return data["photos"][0]["src"]["large"]
    return None

def download_image(image_url):
    response = requests.get(image_url)
    return BytesIO(response.content)

def add_slide_with_image(prs, title, content, image_url, font_name, color_scheme):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    # Title styling
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.name = font_name
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    # Content styling
    content_placeholder = slide.shapes.placeholders[1]
    content_placeholder.text = content
    for p in content_placeholder.text_frame.paragraphs:
        p.font.size = Pt(22)
        p.font.name = font_name
        p.font.color.rgb = RGBColor(*color_scheme['text'])
        p.alignment = PP_ALIGN.LEFT
    if image_url:
        img_stream = download_image(image_url)
        slide.shapes.add_picture(img_stream, Inches(5), Inches(1), Inches(4.5), Inches(3))
    return slide

def apply_slide_transition(slide):
    # Not all pptx libraries support transitions, but we can set a tag for later editing
    slide._element.set('transition', 'fade')

def create_presentation(topics, structure, font_name, color_scheme, branding=None):
    # Ensure color_scheme values are tuples for pptx compatibility
    def ensure_tuple(val):
        if isinstance(val, (list, tuple)) and len(val) == 3:
            return tuple(val)
        return (30, 30, 30)
    color_scheme = {
        'text': ensure_tuple(color_scheme.get('text', (30, 30, 30))),
        'accent': ensure_tuple(color_scheme.get('accent', (79, 140, 255)))
    }
    prs = Presentation()
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = structure['title']
    slide.placeholders[1].text = structure.get('subtitle', '')
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.placeholders[1].text_frame.paragraphs[0].font.size = Pt(28)
    slide.placeholders[1].text_frame.paragraphs[0].font.name = font_name
    if branding and branding.get('logo_url'):
        try:
            logo_stream = download_image(branding['logo_url'])
            slide.shapes.add_picture(logo_stream, Inches(8), Inches(0.2), Inches(1.2), Inches(1.2))
        except Exception:
            pass
    # Section Slides
    if structure.get('sections'):
        for section in structure['sections']:
            add_section_slide(prs, section['title'], font_name, color_scheme)
            for slide_info in section['slides']:
                slide_type = slide_info.get('type', 'bullet')
                slide_title = slide_info.get('title', 'Untitled')
                if slide_type == 'bullet':
                    bullets = openrouter_chat([
                        {"role": "user", "content": f"Generate 3-5 concise, professional bullet points for: {slide_title}"}
                    ])
                    bullets = [b.strip('-• ') for b in bullets.split('\n') if b.strip()]
                    add_bullet_slide(prs, slide_title, bullets, font_name, color_scheme)
                elif slide_type == 'image':
                    content = openrouter_chat([
                        {"role": "user", "content": f"Generate a short, clear summary for: {slide_title}"}
                    ])
                    image_url = search_image(slide_title)
                    add_image_slide(prs, slide_title, content, image_url, font_name, color_scheme)
                elif slide_type == 'chart':
                    chart_prompt = slide_info.get('chart_prompt', f"Generate 3 categories and values for a bar chart about {slide_title}")
                    chart_data = openrouter_chat([
                        {"role": "user", "content": chart_prompt}
                    ])
                    categories, values = [], []
                    for line in chart_data.split('\n'):
                        if ':' in line:
                            cat, val = line.split(':', 1)
                            categories.append(cat.strip())
                            try:
                                values.append(float(val.strip()))
                            except Exception:
                                values.append(1)
                    if categories and values:
                        add_chart_slide(prs, slide_title, slide_title, categories, values, font_name, color_scheme)
    else:
        # Fallback: simple topic list
        for topic in topics:
            bullets = openrouter_chat([
                {"role": "user", "content": f"Generate 3-5 concise, professional bullet points for: {topic}"}
            ])
            bullets = [b.strip('-• ') for b in bullets.split('\n') if b.strip()]
            add_bullet_slide(prs, topic, bullets, font_name, color_scheme)
    # Final slide
    if structure.get('final_slide'):
        add_conclusion_slide(prs, structure['final_slide'].get('title', 'Thank You!'), structure['final_slide'].get('content', 'Questions & Discussion'), font_name, color_scheme)
    tmp = NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(tmp.name)
    return tmp.name

def add_section_slide(prs, title, font_name, color_scheme):
    slide_layout = prs.slide_layouts[2]  # Section Header
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color_scheme['accent'])
    return slide

def add_bullet_slide(prs, title, bullets, font_name, color_scheme):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.name = font_name
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color_scheme['accent'])
    content_placeholder = slide.shapes.placeholders[1]
    content_placeholder.text = ''
    for bullet in bullets:
        p = content_placeholder.text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.name = font_name
        p.font.color.rgb = RGBColor(*color_scheme['text'])
        p.level = 0
    return slide

def add_image_slide(prs, title, content, image_url, font_name, color_scheme):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.name = font_name
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color_scheme['accent'])
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(3)
    if image_url:
        img_stream = download_image(image_url)
        slide.shapes.add_picture(img_stream, left, top, width, height)
    textbox = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(3))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = content
    p.font.size = Pt(20)
    p.font.name = font_name
    p.font.color.rgb = RGBColor(*color_scheme['text'])
    return slide

def add_chart_slide(prs, title, chart_title, categories, values, font_name, color_scheme):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.name = font_name
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color_scheme['accent'])
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(chart_title, values)
    x, y, cx, cy = Cm(2), Cm(4), Cm(16), Cm(8)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.tick_labels.font.size = Pt(14)
    chart.category_axis.tick_labels.font.name = font_name
    chart.value_axis.tick_labels.font.size = Pt(14)
    chart.value_axis.tick_labels.font.name = font_name
    chart.chart_title.text_frame.text = chart_title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(20)
    chart.chart_title.text_frame.paragraphs[0].font.name = font_name
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    return slide

def add_conclusion_slide(prs, title, content, font_name, color_scheme):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.name = font_name
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color_scheme['accent'])
    content_placeholder = slide.shapes.placeholders[1]
    content_placeholder.text = content
    for p in content_placeholder.text_frame.paragraphs:
        p.font.size = Pt(20)
        p.font.name = font_name
        p.font.color.rgb = RGBColor(*color_scheme['text'])
    return slide

def generate_outline_with_llm(topic):
    prompt = f"""
You are a world-class presentation designer and AI assistant. Given the topic: '{topic}', generate a logical, professional outline for a PowerPoint presentation. 
Return a JSON object with:
- title: The main title
- subtitle: A short subtitle
- sections: An array of sections, each with:
    - title: Section title
    - slides: array of slides, each with:
        - type: one of 'bullet', 'image', 'chart' (choose best for content)
        - title: slide title
        - (optional) chart_prompt: for chart slides, a prompt for chart data
- final_slide: {{ "title": ..., "content": ... }} for the conclusion/thank you slide
Example output:
{{
  "title": "The Future of AI",
  "subtitle": "Exploring Innovations, Opportunities, and Challenges",
  "sections": [
    {{"title": "Introduction", "slides": [
      {{"type": "bullet", "title": "What is AI?"}},
      {{"type": "image", "title": "AI in Daily Life"}}
    ]}},
    ...
  ],
  "final_slide": {{"title": "Thank You!", "content": "Questions & Discussion"}}
}}
Be concise, logical, and creative. Use a mix of slide types. Output only valid JSON.
"""
    logger.info(f"[DEBUG] LLM outline prompt: {prompt}")
    outline_json = openrouter_chat([
        {"role": "user", "content": prompt}
    ], max_tokens=1024)
    logger.info(f"[DEBUG] LLM outline raw output: {outline_json}")
    import json
    try:
        outline = json.loads(outline_json)
    except Exception as e:
        logger.error(f"[ERROR] Failed to parse LLM outline JSON: {e}\nRaw: {outline_json}")
        # fallback: single section
        outline = {
            "title": topic,
            "subtitle": "",
            "sections": [
                {"title": topic, "slides": [
                    {"type": "bullet", "title": topic}
                ]}
            ],
            "final_slide": {"title": "Thank You!", "content": "Questions & Discussion"}
        }
    return outline

@router.post("/api/presentation/create")
async def create_presentation_api(request: Request):
    data = await request.json()
    topic = data.get('topic')
    structure = data.get('structure')
    font_name = data.get('font', 'Calibri')
    color_scheme = data.get('color_scheme', {'text': (30, 30, 30), 'accent': (79, 140, 255)})
    branding = data.get('branding', None)
    if not topic or not structure:
        raise HTTPException(status_code=400, detail="Missing topic or structure.")
    # Example: structure = { 'title': ..., 'subtitle': ..., 'slides': [...], 'final_slide': {...} }
    topics = structure.get('slides', [])
    pptx_path = create_presentation(topics, structure, font_name, color_scheme, branding)
    return FileResponse(pptx_path, filename=f"{topic.replace(' ', '_')}.pptx", media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

@router.post("/api/presentation/auto")
async def auto_presentation_api(request: Request):
    try:
        data = await request.json()
        logger.info(f"[DEBUG] Incoming /api/presentation/auto request: {data}")
        topic = data.get('topic')
        font_name = data.get('font', 'Calibri')
        color_scheme = data.get('color_scheme', {'text': (30, 30, 30), 'accent': (79, 140, 255)})
        branding = data.get('branding', None)
        if not topic:
            logger.error("[ERROR] Missing topic in request.")
            raise HTTPException(status_code=400, detail="Missing topic.")
        # Ensure color_scheme is serializable and correct type
        def tupleize(d):
            return {k: tuple(v) if isinstance(v, list) else v for k, v in d.items()}
        color_scheme = tupleize(color_scheme)
        structure = generate_outline_with_llm(topic)
        logger.info(f"[DEBUG] Structure generated: {structure}")
        pptx_path = create_presentation([], structure, font_name, color_scheme, branding)
        logger.info(f"[DEBUG] Presentation created at: {pptx_path}")
        return FileResponse(pptx_path, filename=f"{topic.replace(' ', '_')}.pptx", media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        logger.error(f"[ERROR] Exception in /api/presentation/auto: {e}\nTraceback:\n{tb}")
        raise HTTPException(status_code=500, detail=f"Internal Server Error: {e}\n{tb}")
